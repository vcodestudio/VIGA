"""Evaluation on two PPTX files, predicted and ground truth."""
import argparse
import json
import os
from typing import Any, Dict, List, Optional, Tuple

import pptx
from pptx.enum.shapes import MSO_SHAPE_TYPE

from match import find_maximum_matching
from metrics import get_color_similarity, get_position_similarity, get_shape_fill_similarity, get_text_similarity

def viz_scores(scores: Dict[str, Any]) -> Dict[str, float]:
    """Print the scores in a human-readable format.

    Args:
        scores: Dictionary of metric names to score values (list or scalar).

    Returns:
        Dictionary of metric names to final scores (scaled to 0-100).
    """
    final = {}
    for key, values in scores.items():
        if isinstance(values, list):
            valid_values = [v for v in values if v is not None]
            if len(valid_values) == 0:
                print(f"{key:10s}: N/A")
                final[key] = 100.0
                continue
        else:
            valid_values = values
        if isinstance(valid_values, list) or isinstance(valid_values, int):
            score = sum(valid_values) / len(valid_values)
        else:
            score = valid_values
        print(f"{key:10s}: {score*100:.1f}")
        final[key] = score * 100
    return final


def merge_scores(scores: Dict[str, Any]) -> Dict[str, float]:
    """Merge scores by averaging list values.

    Args:
        scores: Dictionary of metric names to score values.

    Returns:
        Dictionary of metric names to averaged scores.
    """
    merged = {}
    for key, values in scores.items():
        if isinstance(values, list) or isinstance(values, int):
            score = sum(values) / len(values)
        else:
            score = values
        merged[key] = score
    return merged


def extract_text(block: Dict[str, Any]) -> Optional[str]:
    """Extract text content from a block."""
    return block["text"]


def extract_color(block: Dict[str, Any]) -> Any:
    """Extract color/fill from a block."""
    return block["color"]


def extract_position(block: Dict[str, Any]) -> Dict[str, Tuple[float, float, float, float]]:
    """Extract bounding box position from a block."""
    return {"bbox": (
        block["position"][0],  # left
        block["position"][1],  # top
        block["position"][0] + block["size"][0],  # right
        block["position"][1] + block["size"][1],  # bottom
    )}


def parse_blocks(slide: Any, size: Tuple[int, int]) -> List[Dict[str, Any]]:
    """Parse the PPTX slide and extract the blocks.

    Args:
        slide: PPTX slide object.
        size: Slide dimensions as (width, height).

    Returns:
        List of block dictionaries with type, text, color, position, and size.
    """
    blocks = []
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            shape_type = "text"
            text = shape.text_frame.text
            color = shape.fill
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            shape_type = "image"
            text = None
            color = shape.image.blob
        else: # skip other shapes for now
            continue

        blocks.append({
            "type": shape_type, "text": text, "color": color, 
            "position": (shape.left/size[0], shape.top/size[1]),
            "size": (shape.width/size[0], shape.height/size[1]),
        })
    return blocks


def block_match_score(
    gen_blocks: List[Dict[str, Any]],
    ref_blocks: List[Dict[str, Any]],
    matching: List[Tuple[int, int]]
) -> float:
    """Calculate the block match score based on matched areas.

    Args:
        gen_blocks: List of generated blocks.
        ref_blocks: List of reference blocks.
        matching: List of (gen_idx, ref_idx) pairs.

    Returns:
        Match score as ratio of matched area to total area.
    """
    gen_indices = [i for i, _ in matching]
    sum_area, matched_area = 0, 0
    for i, gb in enumerate(gen_blocks):
        if i not in gen_indices:
            sum_area += gb["size"][0] * gb["size"][1]
    ref_indices = [j for _, j in matching]
    for j, rb in enumerate(ref_blocks):
        if j not in ref_indices:
            sum_area += rb["size"][0] * rb["size"][1]
    
    for i, j in matching:
        gb, rb = gen_blocks[i], ref_blocks[j]
        area = gb["size"][0] * gb["size"][1] + rb["size"][0] * rb["size"][1]
        matched_area += area
        sum_area += area
    return matched_area / sum_area


def main() -> None:
    """Run the evaluation on generated vs reference PPTX files."""
    # load the first slide from each PPTX file
    gen_prs = pptx.Presentation(args.generated_pptx)
    gen_slide = gen_prs.slides[args.generated_page-1]
    ref_prs = pptx.Presentation(args.reference_pptx)
    ref_slide = ref_prs.slides[args.reference_page-1]

    # block parsing and matching
    gen_blocks = parse_blocks(gen_slide, size=(gen_prs.slide_width, gen_prs.slide_height))
    ref_blocks = parse_blocks(ref_slide, size=(ref_prs.slide_width, ref_prs.slide_height))
    print(f"# of blocks {len(gen_blocks)} generated {len(ref_blocks)} reference")

    consecutive_bonus, window_size = 1.0, 1
    matching, _, _ = find_maximum_matching(gen_blocks, ref_blocks, consecutive_bonus, window_size)
    matched_blocks = [(gen_blocks[i], ref_blocks[j]) for i, j in matching]
    scores_dict = {"match": block_match_score(gen_blocks, ref_blocks, matching)}

    scores_dict.update({"text": [], "color": [], "position": []})
    for gen_block, ref_block in matched_blocks:
        scores_dict["text"].append(get_text_similarity(
            extract_text(gen_block), extract_text(ref_block)
        ))
        scores_dict["color"].append(get_shape_fill_similarity(
            extract_color(gen_block), extract_color(ref_block)
        ))
        try:
            gen_slide_bg = gen_slide.background.fill.fore_color.rgb
            ref_slide_bg = ref_slide.background.fill.fore_color.rgb
            scores_dict["color"].append(get_color_similarity(
                color1=gen_slide_bg, color2=ref_slide_bg
            ))
        except Exception:
            scores_dict["color"].append(gen_slide.background.fill == ref_slide.background.fill)
        scores_dict["position"].append(get_position_similarity(
            extract_position(gen_block), extract_position(ref_block)
        ))

    final_results = viz_scores(scores_dict)
    if args.output_path is None:
        args.output_path = os.path.join(os.path.dirname(args.generated_pptx), "ref_based.txt")
    with open(args.output_path, 'w') as f:
        for key, value in final_results.items():
            f.write(f"{key}: {value:.1f}\n")


def eval_page(gen_prs: Any, gen_page: int, ref_prs: Any, ref_page: int) -> Dict[str, float]:
    """Evaluate a single page comparison between generated and reference presentations.

    Args:
        gen_prs: Generated PPTX presentation object.
        gen_page: Page index in generated presentation.
        ref_prs: Reference PPTX presentation object.
        ref_page: Page index in reference presentation.

    Returns:
        Dictionary of merged scores for this page.
    """
    gen_slide = gen_prs.slides[gen_page]
    ref_slide = ref_prs.slides[ref_page]

    # block parsing and matching
    gen_blocks = parse_blocks(gen_slide, size=(gen_prs.slide_width, gen_prs.slide_height))
    ref_blocks = parse_blocks(ref_slide, size=(ref_prs.slide_width, ref_prs.slide_height))

    consecutive_bonus, window_size = 1.0, 1
    matching, _, _ = find_maximum_matching(gen_blocks, ref_blocks, consecutive_bonus, window_size)
    matched_blocks = [(gen_blocks[i], ref_blocks[j]) for i, j in matching]
    scores_dict = {"match": block_match_score(gen_blocks, ref_blocks, matching)}

    scores_dict.update({"text": [], "color": [], "position": []})
    for gen_block, ref_block in matched_blocks:
        scores_dict["text"].append(get_text_similarity(
            extract_text(gen_block), extract_text(ref_block)
        ))
        scores_dict["color"].append(get_color_similarity(
            extract_color(gen_block), extract_color(ref_block)
        ))
        scores_dict["position"].append(get_position_similarity(
            extract_position(gen_block), extract_position(ref_block)
        ))
    return merge_scores(scores_dict)



if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Evaluate PPTX file.")
    parser.add_argument("--generated_pptx", type=str, required=True,
                        help="Path to the model-generated PPTX file.")
    parser.add_argument("--reference_pptx", type=str, required=True,
                        help="Path to the reference PPTX file.")
    parser.add_argument("--generated_page", type=int, default=1)
    parser.add_argument("--reference_page", type=int, required=True,
                        help="1-indexed page number.")
    
    parser.add_argument("--output_path", type=str, default=None)

    args = parser.parse_args()

    main()
