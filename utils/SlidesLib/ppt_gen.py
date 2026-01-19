"""PowerPoint presentation generation utilities."""
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.table import Table
from pptx.util import Inches as _Inches
from pptx.util import Pt as _Pt

ARROW_ADD = '"""<a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>"""'


class SlideAgent:
    """Agent for programmatic PowerPoint slide generation."""

    def __init__(
        self, slide_width: float = 13.33, slide_height: float = 7.5
    ) -> None:
        """Initialize a new presentation with specified slide dimensions in inches.

        Args:
            slide_width: Width of slides in inches.
            slide_height: Height of slides in inches.
        """
        self.prs = Presentation()
        self.prs.slide_width = self._inches(slide_width)
        self.prs.slide_height = self._inches(slide_height)
        self.slide: Slide | None = None

    def _inches(self, val: float) -> int:
        """Convert value to EMUs (English Metric Units) via Inches.

        Args:
            val: Value in inches.

        Returns:
            Value in EMUs.
        """
        return _Inches(val)

    def _points(self, val: int) -> int:
        """Convert value to EMUs via Points.

        Args:
            val: Value in points.

        Returns:
            Value in EMUs.
        """
        return _Pt(val)

    # ------- Slide APIs -------
    def add_slide(self, layout: int = 0) -> None:
        """Create a new slide with a specific layout.

        Args:
            layout: Index of the slide layout to use.
        """
        slide_layout = self.prs.slide_layouts[layout]
        self.slide = self.prs.slides.add_slide(slide_layout)

    # ------- Text APIs -------
    def add_title(
        self,
        text: str,
        font_size: int = 44,
        font_color: tuple[int, int, int] = (0, 0, 0),
    ) -> None:
        """Add a title to the slide.

        Args:
            text: Title text.
            font_size: Font size in points.
            font_color: RGB color tuple.
        """
        title_shape = self.slide.shapes.title
        title_shape.text = text
        self._format_text(
            title_shape.text_frame, self._points(font_size), RGBColor(*font_color)
        )

    def add_text(
        self,
        text: str,
        top: float,
        left: float,
        width: float,
        height: float,
        font_size: int = 20,
        bold: bool = False,
        color: tuple[int, int, int] = (0, 0, 0),
        background_color: tuple[int, int, int] | None = None,
        auto_size: bool = True,
    ) -> None:
        """Add a text box at a specified location.

        Args:
            text: Text content.
            top: Top position in inches.
            left: Left position in inches.
            width: Width in inches.
            height: Height in inches.
            font_size: Font size in points.
            bold: Whether text is bold.
            color: RGB color tuple for text.
            background_color: Optional RGB color tuple for background.
            auto_size: Whether to auto-size the text box.
        """
        text_box = self.slide.shapes.add_textbox(
            self._inches(left),
            self._inches(top),
            self._inches(width),
            self._inches(height),
        )

        if background_color:
            text_box.fill.solid()
            text_box.fill.fore_color.rgb = RGBColor(*background_color)
        else:
            text_box.fill.background()

        lines = text.split("\n")
        adjusted_height = height * len(lines)
        text_box.height = self._inches(adjusted_height)

        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        if auto_size:
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self._format_paragraph(
            text_frame, text, self._points(font_size), bold, RGBColor(*color)
        )

    def add_bullet_points(
        self,
        bullet_points: list[str],
        top: float,
        left: float,
        width: float,
        height: float,
        font_size: int = 18,
        color: tuple[int, int, int] = (0, 0, 0),
    ) -> None:
        """Add a text box with bullet points.

        Args:
            bullet_points: List of bullet point strings.
            top: Top position in inches.
            left: Left position in inches.
            width: Width in inches.
            height: Height in inches.
            font_size: Font size in points.
            color: RGB color tuple.
        """
        text_box = self.slide.shapes.add_textbox(
            self._inches(left),
            self._inches(top),
            self._inches(width),
            self._inches(height),
        )
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            self._format_text(p, self._points(font_size), RGBColor(*color))
            p.level = bullet_points.index(point)

    # ------- Image APIs -------
    def add_image(
        self,
        image_path: str,
        top: float,
        left: float,
        width: float,
        height: float,
    ) -> None:
        """Add an image at a specified location.

        Args:
            image_path: Path to the image file.
            top: Top position in inches.
            left: Left position in inches.
            width: Width in inches.
            height: Height in inches.
        """
        self.slide.shapes.add_picture(
            image_path,
            self._inches(left),
            self._inches(top),
            self._inches(width),
            self._inches(height),
        )

    def add_image_centered(
        self, image_path: str, image_width: float, image_height: float
    ) -> None:
        """Add an image centered on the slide.

        Args:
            image_path: Path to the image file.
            image_width: Width of the image in inches.
            image_height: Height of the image in inches.
        """
        slide_width = self.prs.slide_width.inches
        slide_height = self.prs.slide_height.inches
        left = (slide_width - image_width) / 2
        top = (slide_height - image_height) / 2
        self.add_image(image_path, top, left, image_width, image_height)

    # ------- Shape APIs -------
    def add_shape(
        self,
        shape_type: str | int,
        top: float,
        left: float,
        width: float,
        height: float,
        fill_color: tuple[int, int, int] | None = None,
    ) -> None:
        """Add a shape to the slide.

        Args:
            shape_type: Shape type as string or MSO_AUTO_SHAPE_TYPE enum.
            top: Top position in inches.
            left: Left position in inches.
            width: Width in inches.
            height: Height in inches.
            fill_color: Optional RGB color tuple for fill.

        Raises:
            ValueError: If shape_type string is not a valid MSO_AUTO_SHAPE_TYPE.
        """
        if isinstance(shape_type, str):
            try:
                shape_type = getattr(MSO_AUTO_SHAPE_TYPE, shape_type.upper())
            except AttributeError:
                raise ValueError(
                    f"Invalid shape type: {shape_type}. "
                    "Must be a valid MSO_AUTO_SHAPE_TYPE."
                )

        shape = self.slide.shapes.add_shape(
            shape_type,
            self._inches(left),
            self._inches(top),
            self._inches(width),
            self._inches(height),
        )

        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)

    def add_straight_arrow(
        self, start_x: float, start_y: float, end_x: float, end_y: float
    ) -> None:
        """Add a straight arrow connector.

        Args:
            start_x: Starting X position.
            start_y: Starting Y position.
            end_x: Ending X position.
            end_y: Ending Y position.
        """
        self.slide.shapes.add_connector(
            "MSO_CONNECTOR.STRAIGHT", start_x, start_y, end_x, end_y
        )

    def add_straight_line(
        self, start_x: float, start_y: float, end_x: float, end_y: float
    ) -> None:
        """Add a straight line connector.

        Args:
            start_x: Starting X position.
            start_y: Starting Y position.
            end_x: Ending X position.
            end_y: Ending Y position.
        """
        self.slide.shapes.add_connector(
            "MSO_CONNECTOR.STRAIGHT", start_x, start_y, end_x, end_y
        )

    # ------- Table APIs -------
    def add_table(
        self,
        rows: int,
        cols: int,
        top: float,
        left: float,
        width: float,
        height: float,
        column_widths: list[float] | None = None,
    ) -> Table:
        """Add a table to the slide.

        Args:
            rows: Number of rows.
            cols: Number of columns.
            top: Top position in inches.
            left: Left position in inches.
            width: Width in inches.
            height: Height in inches.
            column_widths: Optional list of column widths in inches.

        Returns:
            The created Table object.
        """
        table = self.slide.shapes.add_table(rows, cols, left, top, width, height).table
        if column_widths:
            for idx, col_width in enumerate(column_widths):
                table.columns[idx].width = self._inches(col_width)
        return table

    # ------- Helper APIs -------
    def set_background_color(self, color: tuple[int, int, int]) -> None:
        """Set background color for the current slide.

        Args:
            color: RGB color tuple.
        """
        background = self.slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    def duplicate_slide(self, slide_index: int) -> None:
        """Duplicate a slide by index.

        Args:
            slide_index: Index of the slide to duplicate.
        """
        template_slide = self.prs.slides[slide_index]
        new_slide = self.prs.slides.add_slide(template_slide.slide_layout)
        for shape in template_slide.shapes:
            self._copy_shape(shape, new_slide)

    def save_presentation(self, file_name: str) -> None:
        """Save the PowerPoint presentation.

        Args:
            file_name: Output file path.
        """
        self.prs.save(file_name)

    # ------- Internal Helper Methods -------
    def _format_paragraph(
        self,
        text_frame: Any,
        text: str,
        font_size: int,
        bold: bool,
        color: RGBColor,
    ) -> None:
        """Format text within a text frame.

        Args:
            text_frame: The text frame to format.
            text: Text content.
            font_size: Font size in EMUs.
            bold: Whether text is bold.
            color: RGBColor for text.
        """
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color

    def _format_text(
        self, text_frame: Any, font_size: int, font_color: RGBColor
    ) -> None:
        """Format text in a text frame.

        Args:
            text_frame: The text frame to format.
            font_size: Font size in EMUs.
            font_color: RGBColor for text.
        """
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.color.rgb = font_color

    def _copy_shape(self, shape: BaseShape, slide: Slide) -> None:
        """Copy a shape from one slide to another.

        Args:
            shape: The shape to copy.
            slide: The target slide.
        """
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = BytesIO(shape.image.blob)
            slide.shapes.add_picture(
                image, shape.left, shape.top, shape.width, shape.height
            )
        elif shape.has_text_frame:
            new_shape = slide.shapes.add_textbox(
                shape.left, shape.top, shape.width, shape.height
            )
            new_shape.text = shape.text
            self._format_text(
                new_shape.text_frame,
                shape.text_frame.paragraphs[0].font.size,
                shape.text_frame.paragraphs[0].font.color.rgb,
            )
