import os
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import  MSO_AUTO_SIZE


def get_text_format(shape, default_font_size: int = 18):
    """Extract font size and color from a text shape."""
    font_size = None   # Default to None if not found
    color = (0, 0, 0)  # Default to black
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            # Set default values for font size and color
            if run.font.size:
                font_size = run.font.size.pt
            
            if run.font.color:
                if run.font.color.type == 1:  # RGB color
                    color_rgb = run.font.color.rgb
                    if color_rgb:
                        color = (color_rgb[0], color_rgb[1], color_rgb[2])
                elif run.font.color.type == 2:  # Scheme color
                    # Handle the theme color here
                    # You can set a default value or leave it as is
                    color = (0, 0, 0)  # You can change this to any default color

            if font_size is None:
                font_size = default_font_size
            return font_size, color

    # Default values if no runs or paragraphs exist
    return default_font_size, (0, 0, 0)


def get_shape_background_color(shape):
    """Extract background color from a shape."""
    if shape.fill.type == 1:  # Solid fill
        color = shape.fill.fore_color
        if color.type == 1:  # RGB color
            return (color.rgb[0], color.rgb[1], color.rgb[2])
        elif color.type == 2:  # Scheme color
            # Scheme colors refer to theme colors, and an RGB equivalent might not always be directly available
            # We'll just indicate it's a theme color for now
            return None  # Or handle based on your needs, e.g., default RGB color
    return None


def get_slide_background_color(slide):
    """Extract background color from a slide."""
    bg_fill = slide.background.fill
    if bg_fill.type == 1:  # Solid fill
        color = bg_fill.fore_color
        if color.type == 1:  # RGB color
            return (color.rgb[0], color.rgb[1], color.rgb[2])
        elif color.type == 2:  # Scheme color
            # Handle scheme color by assigning a default RGB color or returning None
            return None  # or return (255, 255, 255) as a default white color if preferred
    return None


def copy_shape_code(f, shape, parent_name):
    if shape.has_text_frame:
        # Create a text box instead of a general auto shape
        f.write(f"    shape = {parent_name}.shapes.add_textbox(Inches({shape.left.inches}), Inches({shape.top.inches}), Inches({shape.width.inches}), Inches({shape.height.inches}))\n")
    else:
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_type_str = f"MSO_AUTO_SHAPE_TYPE.{shape.auto_shape_type}"
        else:
            shape_type_str = f"MSO_SHAPE_TYPE.{shape.shape_type}"
        shape_type_str = shape_type_str.split(" (")[0]
        f.write(f"    shape = {parent_name}.shapes.add_shape({shape_type_str}, Inches({shape.left.inches}), Inches({shape.top.inches}), Inches({shape.width.inches}), Inches({shape.height.inches}))\n")
    
    # Apply fill color if applicable
    if hasattr(shape, 'fill') and shape.fill:
        if shape.fill.type == 1:  # Solid fill
            f.write(f"    shape.fill.solid()\n")
            if shape.fill.fore_color.type == 1:  # RGB color
                f.write(f"    shape.fill.fore_color.rgb = RGBColor({shape.fill.fore_color.rgb[0]}, {shape.fill.fore_color.rgb[1]}, {shape.fill.fore_color.rgb[2]})\n")
            elif shape.fill.fore_color.type == 2:  # Scheme color
                scheme_color_str = f"MSO_THEME_COLOR.{shape.fill.fore_color.theme_color}"
                scheme_color_str = scheme_color_str.split(" (")[0]
                f.write(f"    shape.fill.fore_color.theme_color = {scheme_color_str}\n")
            else:
                # Handle other non-RGB color cases
                f.write(f"    # Non-RGB color detected. Manual intervention may be required.\n")
                f.write(f"    # shape.fill.fore_color is of type {shape.fill.fore_color.type}\n")
                # Optionally, apply a default color or handle the case based on your needs
                f.write(f"    shape.fill.fore_color.rgb = RGBColor(128, 128, 128)  # Placeholder for non-RGB color\n")
    
    
    # Apply line color and width if applicable
    if hasattr(shape, 'line') and shape.line and shape.line.color.type == 1:  # RGB color
        f.write(f"    shape.line.color.rgb = RGBColor({shape.line.color.rgb[0]}, {shape.line.color.rgb[1]}, {shape.line.color.rgb[2]})\n")
        f.write(f"    shape.line.width = {shape.line.width}\n")
    elif hasattr(shape, 'line') and shape.line and shape.line.color.type == 2:  # Scheme color
        scheme_color_str = f"MSO_THEME_COLOR.{shape.line.color.theme_color}"
        scheme_color_str = scheme_color_str.split(" (")[0]
        f.write(f"    shape.line.color.theme_color = {scheme_color_str}\n")
    
    if shape.has_text_frame:
        # Clear existing text
        f.write(f"    shape.text_frame.clear()\n")
        # Apply rotation if any
        if shape.rotation:
            f.write(f"    shape.rotation = {shape.rotation}\n")
        # Apply text wrapping and auto-sizing
         # Apply text wrapping and auto-sizing
        f.write(f"    shape.text_frame.word_wrap = {shape.text_frame.word_wrap}\n")
        if shape.text_frame.auto_size == MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT:
            f.write(f"    shape.text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT\n")
        elif shape.text_frame.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE:
            f.write(f"    shape.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE\n")
        else:
            f.write(f"    shape.text_frame.auto_size = None\n")
            
        # Apply text margins
        f.write(f"    shape.text_frame.margin_left = {shape.text_frame.margin_left}\n")
        f.write(f"    shape.text_frame.margin_right = {shape.text_frame.margin_right}\n")
        f.write(f"    shape.text_frame.margin_top = {shape.text_frame.margin_top}\n")
        f.write(f"    shape.text_frame.margin_bottom = {shape.text_frame.margin_bottom}\n")
        
        for p_idx, p in enumerate(shape.text_frame.paragraphs):
            # Add paragraphs
            if p_idx > 0:
                f.write(f"    p = shape.text_frame.add_paragraph()\n")
            else:
                f.write(f"    p = shape.text_frame.paragraphs[0]\n")

            # Set paragraph alignment and spacing
            if p.alignment:
                alignment_str = f"PP_ALIGN.{p.alignment}"
                alignment_str = alignment_str.split(' (')[0]
                f.write(f"    p.alignment = {alignment_str}\n")
            f.write(f"    p.space_after = {p.space_after}\n")
            f.write(f"    p.space_before = {p.space_before}\n")

            # Clear paragraph text before adding runs
            f.write(f"    p.clear()\n")

            for run in p.runs:
                # Add run and set text and formatting
                f.write(f"    run = p.add_run()\n")
                f.write(f"    run.text = {run.text!r}\n")
                if run.font.size:
                    f.write(f"    run.font.size = Pt({run.font.size.pt})\n")
                f.write(f"    run.font.bold = {run.font.bold}\n")
                f.write(f"    run.font.italic = {run.font.italic}\n")
                f.write(f"    run.font.underline = {run.font.underline}\n")
                if run.font.color.type == 1:  # RGB color
                    f.write(f"    run.font.color.rgb = RGBColor({run.font.color.rgb[0]}, {run.font.color.rgb[1]}, {run.font.color.rgb[2]})\n")
                f.write(f"    run.font.name = {run.font.name!r}\n")
    
# Note: You need to save the high-level APIs (add_title, add_text, etc.) in a file named 'high_level_apis.py'
# and place it in the same directory as the generated script or adjust the import path accordingly.

# %% Main Function

PPTX_IMPORTS = [
    "from pptx import Presentation\n",
    "from pptx.util import Inches, Pt\n",
    "from pptx.dml.color import RGBColor\n",
    "from pptx.enum.text import MSO_AUTO_SIZE\n",
    "from pptx.enum.shapes import MSO_SHAPE_TYPE\n\n"
]

def format_coords(coords: list[float], n: int = 2) -> str:
    return '[' + ', '.join([str(round(c, n)) for c in coords]) + ']'

def get_slide_code(slide, output_dir: str, slide_shape: tuple[int, int]):
    """Generate canonical code to reproduce one slide."""
    output_py = os.path.join(output_dir, 'code_library.py')

    contents = PPTX_IMPORTS + ["from library import *\n\n"] # import library functions
    contents.append("prs = Presentation()\n")
    # set slide size
    slide_width, slide_height = slide_shape
    contents.append(f"prs.slide_width = {slide_width}\n")
    contents.append(f"prs.slide_height = {slide_height}\n\n")
    # add slide
    contents.append(f"slide_layout = prs.slide_layouts[5]  # blank layout\n")  # zora: always use blank layout?
    contents.append(f"slide = prs.slides.add_slide(slide_layout)\n\n")
    # handle background color
    bg_color = get_slide_background_color(slide)
    if bg_color:
        contents.append(f"set_background_color(slide, ({bg_color[0]}, {bg_color[1]}, {bg_color[2]}))\n\n")
    
    # process shapes
    image_counter = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Save the image
            image_path = os.path.join(output_dir, f'image_{image_counter}.jpg')
            with open(image_path, 'wb') as img_file:
                img_file.write(shape.image.blob)
            # Write code to add image using high-level API
            coords = [shape.left.inches, shape.top.inches, shape.width.inches, shape.height.inches]
            coords = format_coords(coords)
            contents.append(f"add_image(slide, 'media/{os.path.basename(image_path)}', {coords})\n")
            image_counter += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape.placeholder_format.type == 1:
            # Title placeholder
            text = shape.text_frame.text
            font_size, font_color = get_text_format(shape)
            bg_color = get_shape_background_color(shape)
            contents.append(f"add_title(slide, {text!r}, font_size={int(font_size)}, font_color={font_color}, background_color={bg_color})\n")
        elif shape.has_text_frame:
            # determine if it's bullet points or regular text
            texts = [p.text for p in shape.text_frame.paragraphs]
            is_bullet = any(p.level > 0 for p in shape.text_frame.paragraphs)
            coords = [shape.left.inches, shape.top.inches, shape.width.inches, shape.height.inches]
            coords = format_coords(coords)
            font_size, font_color = get_text_format(shape)
            bg_color = get_shape_background_color(shape)
            if is_bullet:
                contents.append(f"add_bullet_points(slide, {texts}, {coords}, font_size={int(font_size)}, color={font_color}, background_color={bg_color})\n")
            else:
                text = "\n".join(texts)
                if len(text) > 0:
                    contents.append(f"add_text(slide, {text!r}, {coords}, font_size={int(font_size)}, color={font_color}, background_color={bg_color})\n")
        else:
            pass

    contents.append("prs.save('output.pptx')")

    with open(output_py, 'w') as f:
        f.write(''.join(contents))



def generate_code_from_pptx(input_pptx: str, output_path: str):
    """Generate canonical code to reproduce all slides using high-level APIs."""
    # Load the input presentation
    prs = Presentation(input_pptx)
    for i, slide in enumerate(prs.slides):
        output_dir = f"{output_path}/slide_{i+1}/"
        os.makedirs(output_dir, exist_ok=True)

        slide_shape = (prs.slide_width, prs.slide_height)
        get_slide_code(slide, output_dir, slide_shape)



if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Generate Groundtruth Code to reproduce the slide.')
    parser.add_argument("--slides_path", type=str, required=True, help="Path to the pptx file.")
    parser.add_argument("--output_dir", type=str, required=True, help="Path to output the example.")
    args = parser.parse_args()

    generate_code_from_pptx(
        input_pptx=args.slides_path,
        output_path=args.output_dir
    )