"""Reference-free slide evaluation using VLM scoring."""
import argparse
import base64
import json
import os
import sys

import pptx
from openai import OpenAI
from PIL import Image

sys.path.append(os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))
from utils.common import get_model_info

INSTRUCTION = """Please evaluate the slide based on the following criteria:
{}

Give an integer score between 0 - {}, higher scores means the criteria is better met.
First, respond with a score; Then, provide your justification for the score in natural language sentences. Your response should look like this: '4. The slide has concise texts summarized in bullet points.'
Only evaluate the slide based on the specified criteria, and no other aspects. Give scores across the full spectrum (1-5) instead of only good ones (3-5).
""" # + "Try give more polarized scores (e.g., 1 or 5) instead of mediocore ones (e.g., 2 or 3)."


METRIC_DICT = {
    "text": {"criteria": "The title should be simple and clear to indicate the main point. For main content, avoid too many texts and keep words concise. Use a consistent and readable font size, style, and color.", "scale": 5},
    "image": {"criteria": "Use high-quality images with a reasonable proportion. Do not penalize the slide if no image is involved.", "scale": 5},
    "layout": {"criteria": "Elements should be aligned, do not overlap, and have sufficient margins to each other. All elements should not exceed the page.", "scale": 5},
    "color": {"criteria": "Use high-contrast color especially between the text and the background. Avoid using high-glaring colors.", "scale": 5},
}

def encode_image(image_path):
    with open(image_path, 'rb') as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

def evaluate_slide(prompt: str, image_url: str) -> str:
    messages = [{
        "role": "user",
        "content": [
            {"type": "text", "text": prompt},
            {"type": "image_url", "image_url": {"url": image_url}},
        ],
    }]
    client = OpenAI(api_key=get_model_info(args.model_name)["api_key"])
    response = client.chat.completions.create(
            model=args.model_name,
            messages=messages,
            max_tokens=args.max_tokens,
            temperature=0.0,
    )
    response = response.choices[0].message.content.strip()
    try:
        score, justification = response.split(".", 1)
        score = float(score.strip())
    except:
        score, justification = 0.0, response
    return {
        "score": score,
        "justification": justification.strip(),
    }

def parse_text(path, page: int = 0):
    slide = pptx.Presentation(path).slides[page]
    text_list = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_list.append(shape.text_frame.text)
    text_list = [t for t in text_list if t.strip()]
    return text_list


def main():
    # convert png to jpg
    if args.image_path.endswith(".png"):
        jpg_image_path = args.image_path.replace('.png', '.jpg')
        if os.path.exists(jpg_image_path):
            print("JPG image already exists! Skip conversion.")
        png_image = Image.open(args.image_path)
        rgb_image = png_image.convert("RGB")
        rgb_image.save(jpg_image_path, format="JPEG", quality=95)
        print("Conversion successful! Saved as ", jpg_image_path)
    else:
        jpg_image_path = args.image_path

    image_base64 = encode_image(jpg_image_path)
    image_url = f"data:image/jpeg;base64,{image_base64}"

    results_dict = {}
    for metric, metric_dict in METRIC_DICT.items():
        if metric not in args.metric_to_use: continue
        prompt = INSTRUCTION.format(metric_dict["criteria"], metric_dict["scale"])
        if metric == "text":
            pptx_path = jpg_image_path.replace(".jpg", ".pptx")
            slide_texts = parse_text(pptx_path)
            prompt += '\n\nTexts in the slide are: \n' + '\n'.join(slide_texts)
        results_dict[metric] = evaluate_slide(prompt, image_url)

    for metric, result in results_dict.items():
        print(f"{metric}: {result}\n")
    print("Total: ", sum([v["score"] for v in results_dict.values()]))
    
    if args.response_path is None:
        response_path = jpg_image_path.replace(".jpg", "_ref_free.json")
    else:
        response_path = args.response_path
    print(f"Saving response to {response_path}")
    json.dump(results_dict, open(response_path, "w"))



if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="PPTX reference-free evaluation.")
    parser.add_argument("--image_path", type=str, required=True,
                        help="Path to the image of model-generated PPTX file.")
    parser.add_argument("--response_path", type=str, default=None,
                        help="Path to save the model-evaluation justifications.")
    parser.add_argument("--metric_to_use", type=str, nargs='+',
                        default=["text", "image", "layout", "color"],
                        help="Metrics to evaluate.")
    parser.add_argument("--model_name", type=str, default="gpt-4o", help="Model name to use.")
    parser.add_argument("--max_tokens", type=int, default=100, help="Max tokens to generate.")

    args = parser.parse_args()
    main()
