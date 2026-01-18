import os
import base64
import argparse
from pptx import Presentation

import openai
openai.api_key = os.environ["OPENAI_API_KEY"]
client = openai.OpenAI()


# INSTRUCTION = """Your task is to write a clear instruction on how to create the provided slide.
# Make sure you include all necessary information about content, formatting, and other important aspects.
# Do not over-specify things. Please phrase the instruction using the fewest words possible."""

INSTRUCTION = """Your task is to write a high-level instruction to create the provided slide.
Do not over-specify things. Please phrase the instruction using the fewest words possible, ideally in one sentence.
"""


def encode_image(image_path):
  with open(image_path, 'rb') as image_file:
    return base64.b64encode(image_file.read()).decode('utf-8')


def main(pptx_path, output_path):
    prs = Presentation(pptx_path)

    # prepare system prompt and few-shot examples
    messages = [{"role": "system", "content": INSTRUCTION}]
    for i, slide in enumerate(prs.slides):
        if i > (args.num_seeds - 1): break
        ist_path = os.path.join(output_path, f'slide_{i+1}', f'instruction_high_level.txt')
        # ist_path = os.path.join(output_path, f'slide_{i+1}', f'instruction_human.txt')
        instruction = open(ist_path).read()
        img_path = os.path.join(output_path, f'slide_{i+1}', f'slide.png')
        image_url = f"data:image/jpeg;base64,{encode_image(img_path)}"
        messages += [{
            "role": "user",
            "content": [
                {"type": "text", "text": "Write a clear instruction on how to create the following slide."},
                {"type": "image_url", "image_url": {"url": image_url}},
            ],
        }]
        messages += [{
            "role": "assistant",
            "content": [
                {"type": "text", "text": instruction},
            ],
        }]
    
    # process each slide (in image format) and save the instruction to generated this slide.
    for i, slide in enumerate(prs.slides):
        img_dir = os.path.join(output_path, f'slide_{i+1}')
        img_path = os.path.join(img_dir, f'slide.png')
        image_base64 = encode_image(img_path)
        image_url = f"data:image/jpeg;base64,{image_base64}"
        add_messages = []
        add_messages += [{
            "role": "user",
            "content": [
                {"type": "text", "text": f'Example {i}:\n' + instruction},
                {"type": "image_url", "image_url": {"url": image_url}},
            ],
        }]
        add_messages += [{
            "role": "user",
            "content": [
                {"type": "text", "text": "Now write a clear instruction on how to create the following slide."},
                {"type": "image_url", "image_url": {"url": image_url}},
            ],
        }]
        response = client.chat.completions.create(
            model=args.model_name,
            messages=messages + add_messages,
            max_tokens=args.max_tokens,
            temperature=0.0,
        )
        print(response)
        instruction = response.choices[0].message.content
        instruction_path = os.path.join(img_dir, f'instruction_high_level_model.txt')
        # instruction_path = os.path.join(img_dir, f'instruction_model.txt')
        with open(instruction_path, 'w') as fw:
            fw.write(instruction)
        print(f"Slide {i+1} instruction generated.")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Generate instruction')
    parser.add_argument("--pptx_path", type=str, required=True, help="Path to the pptx-formed slide file.")
    parser.add_argument("--model_name", type=str, default="gpt-4o-mini", help="Model name to use.")
    parser.add_argument("--max_tokens", type=int, default=300, help="Max tokens to generate.")
    parser.add_argument("--num_seeds", type=int, default=3, help="Number of seed instructions.")
    parser.add_argument("--output_path", type=str, required=True)
    args = parser.parse_args()

    main(args.pptx_path, args.output_path)
