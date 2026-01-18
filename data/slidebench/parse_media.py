import os
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def collect_unique_images(slide) -> list:
    """Find all non-duplicate images in a slide."""
    def iter_picture_shapes(slide):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape

    ibytes_set = set()
    for picture in iter_picture_shapes(slide):
        image_bytes = picture.image.blob  # get image "file" contents
        if image_bytes not in ibytes_set:
            ibytes_set.add(image_bytes)
    return list(ibytes_set)


def save_images_in_slide(slide, output_dir: str, slide_number: int) -> None:
    """Save all unique images in a slide to the output directory."""
    image_bytes_list = collect_unique_images(slide)

    # Add slide number for a full presentation.
    slide_dir = os.path.join(output_dir, f'slide_{slide_number}')
    os.makedirs(slide_dir, exist_ok=True)
    image_dir = os.path.join(slide_dir, 'media')   
    os.makedirs(image_dir, exist_ok=True)
    for i, image_bytes in enumerate(image_bytes_list):
        image_name = os.path.join(image_dir, f'image_{i}.jpg')
        with open(image_name, 'wb') as f:
            f.write(image_bytes)



def main():
    prs = Presentation(args.slides_path)  # load presentation
    
    if args.slide_index is not None:  # process a single slide
        slide_index = args.slide_index - 1  # convert 1-based index to 0-based
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise ValueError(f"Slide index {args.slide_index} is out of range.")
        slides_to_process = [(prs.slides[slide_index], args.slide_index)]
    else:  # process all slides
        slides_to_process = [(slide, i + 1) for i, slide in enumerate(prs.slides)]
    
    # save the media files
    os.makedirs(args.output_dir, exist_ok=True)
    
    for slide, slide_number in slides_to_process:
        save_images_in_slide(slide, args.output_dir, slide_number)
        # save_media_in_slide(slide, args.output_dir, slide_number, "video")
        # save_media_in_slide(slide, args.output_dir, slide_number, "audio")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description='Parse Media from PPTX')
    parser.add_argument("--slides_path", type=str, required=True, help="Path to the pptx file.")
    parser.add_argument("--output_dir", type=str, required=True, help="Path to output the example.")
    parser.add_argument("--slide_index", type=int, required=False, 
                        help="Slide index to process (1-based). If not specified, process all slides.")
    args = parser.parse_args()

    main()
