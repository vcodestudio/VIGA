from pptx import Presentation
from pptx.util import Inches as _Inches, Pt as _Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from io import BytesIO

ARROW_ADD = '"""<a:tailEnd type="arrow" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>"""'

class SlideAgent:
    def __init__(self, slide_width=13.33, slide_height=7.5):
        """Initialize a new presentation with specified slide dimensions in inches."""
        self.prs = Presentation()
        self.prs.slide_width = self._inches(slide_width)
        self.prs.slide_height = self._inches(slide_height)
        self.slide = None

    def _inches(self, val):
        """Helper method to convert to Inches."""
        return _Inches(val)

    def _points(self, val):
        """Helper method to convert to Points."""
        return _Pt(val)

    # ------- Slide APIs -------
    def add_slide(self, layout=0):
        """Create a new slide with a specific layout."""
        slide_layout = self.prs.slide_layouts[layout]
        self.slide = self.prs.slides.add_slide(slide_layout)

    # ------- Text APIs -------
    def add_title(self, text, font_size=44, font_color=(0, 0, 0)):
        """Add a title to the slide with a custom font size (in points) and font color (RGB tuple)."""
        title_shape = self.slide.shapes.title
        title_shape.text = text
        self._format_text(title_shape.text_frame, self._points(font_size), RGBColor(*font_color))

    def add_text(self, text, top, left, width, height, font_size=20, bold=False, color=(0, 0, 0), background_color=None, auto_size=True):
        """Add a text box at a specified location with custom text settings and optional background color."""
        # Create the text box shape
        text_box = self.slide.shapes.add_textbox(self._inches(left), self._inches(top), self._inches(width), self._inches(height))

        # Set background color if provided
        if background_color:
            text_box.fill.solid()
            text_box.fill.fore_color.rgb = RGBColor(*background_color)
        else:
            text_box.fill.background()  # No fill if no color is specified

        # Handle line breaks and adjust height
        lines = text.split("\n")
        adjusted_height = height * len(lines)  # Adjust height based on the number of lines
        text_box.height = self._inches(adjusted_height)

        # Set text and format it
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        if auto_size:
            text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT  # Automatically fit the text box to the text
        self._format_paragraph(text_frame, text, self._points(font_size), bold, RGBColor(*color))

    def add_bullet_points(self, bullet_points, top, left, width, height, font_size=18, color=(0, 0, 0)):
        """Add a text box with bullet points."""
        text_box = self.slide.shapes.add_textbox(self._inches(left), self._inches(top), self._inches(width), self._inches(height))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        for point in bullet_points:
            p = text_frame.add_paragraph()
            p.text = point
            self._format_text(p, self._points(font_size), RGBColor(*color))
            p.level = bullet_points.index(point)

    # ------- Image APIs -------
    def add_image(self, image_path, top, left, width, height):
        """Add an image at a specified location."""
        self.slide.shapes.add_picture(image_path, self._inches(left), self._inches(top), self._inches(width), self._inches(height))

    def add_image_centered(self, image_path, image_width, image_height):
        """Add an image centered on the slide."""
        slide_width = self.prs.slide_width.inches
        slide_height = self.prs.slide_height.inches
        left = (slide_width - image_width) / 2
        top = (slide_height - image_height) / 2
        self.add_image(image_path, top, left, image_width, image_height)
        
    # ------- Shape APIs -------
    def add_shape(self, shape_type, top, left, width, height, fill_color=None):
        """Add a shape to the slide, supporting MSO_AUTO_SHAPE_TYPE."""
        if isinstance(shape_type, str):
            # Check if the shape type is a valid string, otherwise raise an error
            try:
                shape_type = getattr(MSO_AUTO_SHAPE_TYPE, shape_type.upper())
            except AttributeError:
                raise ValueError(f"Invalid shape type: {shape_type}. Must be a valid MSO_AUTO_SHAPE_TYPE.")
        
        # Now create the shape with the validated or passed enum type
        shape = self.slide.shapes.add_shape(shape_type, self._inches(left), self._inches(top), self._inches(width), self._inches(height))
        
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*fill_color)
            
    def add_straight_arrow(self, start_x, start_y, end_x, end_y):
        connector = self.slide.shapes.add_connector("MSO_CONNECTOR.STRAIGHT", start_x, start_y, end_x, end_y)
        
    
    def add_straight_line(self, start_x, start_y, end_x, end_y):
        connector = self.slide.shapes.add_connector("MSO_CONNECTOR.STRAIGHT", start_x, start_y, end_x, end_y)
        line_elem = connector.line._get_or_add_ln()
        line_elem.append(parse_xml({ARROW_ADD}))

    # ------- Table APIs -------
    def add_table(self, rows, cols, top, left, width, height, column_widths=None):
        """Add a table to the slide."""
        table = self.slide.shapes.add_table(rows, cols, left, top, width, height).table
        if column_widths:
            for idx, col_width in enumerate(column_widths):
                table.columns[idx].width = Inches(col_width)
        return table

    # ------- Helper APIs -------
    def set_background_color(self, color):
        """Set background color for the current slide."""
        background = self.slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def duplicate_slide(self, slide_index):
        """Duplicate a slide by index."""
        template_slide = self.prs.slides[slide_index]
        new_slide = self.prs.slides.add_slide(template_slide.slide_layout)
        for shape in template_slide.shapes:
            self._copy_shape(shape, new_slide)

    def save_presentation(self, file_name):
        """Save the PowerPoint presentation."""
        self.prs.save(file_name)

    # ------- Internal Helper Methods -------
    def _format_paragraph(self, text_frame, text, font_size, bold, color):
        """Helper function to format text within a text frame."""
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = font_size
        p.font.bold = bold
        p.font.color.rgb = color

    def _format_text(self, text_frame, font_size, font_color):
        """Helper function to format text in a text frame."""
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = font_size
            paragraph.font.color.rgb = font_color

    def _copy_shape(self, shape, slide):
        """Copy a shape from one slide to another."""
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image = BytesIO(shape.image.blob)
            slide.shapes.add_picture(image, shape.left, shape.top, shape.width, shape.height)
        elif shape.has_text_frame:
            new_shape = slide.shapes.add_textbox(shape.left, shape.top, shape.width, shape.height)
            new_shape.text = shape.text
            self._format_text(new_shape.text_frame, shape.text_frame.paragraphs[0].font.size, shape.text_frame.paragraphs[0].font.color.rgb)
        



    

 